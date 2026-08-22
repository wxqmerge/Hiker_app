from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1 - Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
title.text = "Hiker App User Guide"
subtitle = slide1.placeholders[1]
subtitle.text = "Next Hike Banner - User Use Case\nGreen Banner Quick Reference"
# Style
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True

# Slide 2 - What is the Green Banner
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "What is the Green Banner?"
content2 = slide2.placeholders[1]
content2.text = """The Next Hike Banner appears at the top of every page when you have an upcoming hike scheduled.

• Shows your NEXT scheduled hike
• Displays key information at a glance
• Auto-updates based on your schedule
• Designed for phone use - large, readable text

Perfect for hikers checking conditions day-of"""
# Style
title2.text_frame.paragraphs[0].font.size = Pt(32)

# Slide 3 - Banner Components
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "Banner Components"
content3 = slide3.placeholders[1]
content3.text = """TRAIL NAME
Large, clear text showing hike name

DATE TILE
Day of week + date (e.g., "Wed 24 Sep")

DIFFICULTY & EARLY START
Color-coded difficulty rating
Early start time if applicable

WEATHER INFO
Temperature (high/low)
Rain probability

TIDE INFO
Low tide time and height (for coastal hikes)

TRAIL CONDITIONS
Quick status indicators"""

# Slide 4 - User Use Case A
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Use Case A: Phone Hikers Day-of"
content4 = slide4.placeholders[1]
content4.text = """Target User: Hikers on phone morning-of hike
Typical users: Older adults, need large readable text

What you need to know:
✓ What hike is next
✓ What time to meet
✓ Weather conditions
✓ Tide conditions (if coastal)
✓ Trail difficulty

No scrolling needed - all info visible at once"""

# Slide 5 - How to Read Weather
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Reading Weather Information"
content5 = slide5.placeholders[1]
content5.text = """Temperature
• Shows high/low range for the day
• Updated automatically from weather service
• Example: 18°C / 12°C

Rain Probability
• % chance of rain
• Helps decide what to bring
• Updated in real-time

Note: Weather data is fetched automatically when banner loads.
If no weather data, shows "N/A" - trail may not have coordinates"""

# Slide 6 - Tide Information
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Tide Information"
content6 = slide6.placeholders[1]
content6.text = """For coastal hikes with tide stations:

Low Tide Time
• Shows exact time of low tide
• Critical for beach access

Low Tide Height
• Height in meters/feet
• Shows if beach is accessible

Important: Some trails have tide data WITHOUT weather coords.
Banner shows tide even if weather is N/A.

Example: Low: 10:42 AM, 0.3m"""

# Slide 7 - Difficulty & Timing
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Difficulty & Timing"
content7 = slide7.placeholders[1]
content7.text = """Difficulty Rating
• Color-coded (Green/Yellow/Red)
• Consistent with trail database

Early Start Time
• Shows if hike starts early
• Example: "Early Start 7:30 AM"
• Important for sunrise hikes

Date Information
• Day of week and date
• Next hike only (not entire month)"""

# Slide 8 - When Banner Appears
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
title8.text = "When Does Banner Appear?"
content8 = slide8.placeholders[1]
content8.text = """Banner shows when:

✓ You have a hike scheduled in the future
✓ Schedule is loaded
✓ Next hike is determined

Banner hides when:
✗ No upcoming hikes scheduled
✗ Schedule is empty

The banner updates automatically:
• When you change schedule
• When page loads
• When weather data refreshes"""

# Slide 9 - Mobile Optimization
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
title9.text = "Mobile Optimization"
content9 = slide9.placeholders[1]
content9.text = """Designed for phone use:

✓ Large text for older adults
✓ Single screen view - no scrolling
✓ Touch-friendly layout
✓ Responsive design
✓ High contrast green background
✓ Clear typography

Viewport optimized:
• Initial scale: 0.9 for better fit
• Banner width adapts to screen
• Information prioritized"""

# Slide 10 - Tips & Best Practices
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title10 = slide10.shapes.title
title10.text = "Tips for Hikers"
content10 = slide10.placeholders[1]
content10.text = """Before Hike:
1. Check banner morning-of for weather
2. Note tide times for coastal hikes
3. Confirm early start time
4. Pack accordingly

If banner shows N/A:
• Weather data may be unavailable
• Check trail has coordinates in database
• Contact admin if persistent

Banner updates automatically - refresh page if needed"""

# Save
prs.save('D:\\hiker\\Hiker_App_User_Guide_Next_Hike_Banner.pptx')
print("Presentation created successfully!")
